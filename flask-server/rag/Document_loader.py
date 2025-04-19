from langchain_community.document_loaders import DirectoryLoader
from rag.RAG import DocumentLoaderFactory
from langchain.document_loaders.base import BaseLoader
from langchain_community.document_loaders import (
    PyPDFLoader,
    UnstructuredPowerPointLoader,
    UnstructuredWordDocumentLoader,
    TextLoader,
    UnstructuredFileLoader,
    CSVLoader,
    JSONLoader,
    UnstructuredHTMLLoader,
    UnstructuredMarkdownLoader,
    BSHTMLLoader,
    UnstructuredExcelLoader,
    PythonLoader,
    # UnstructuredImageLoader,
    # UnstructuredVideoLoader,
    # UnstructuredAudioLoader,
)
from pptx import Presentation
from langchain.docstore.document import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Tuple, Any
import os
import logging
import datetime


class PDFDirectoryLoaderFactory(DocumentLoaderFactory):
    def create_loader(self, folder_path: str) -> BaseLoader:
        return DirectoryLoader(folder_path, glob="*.pdf", loader_cls=PyPDFLoader)


class MultiDocumentDirectoryLoaderFactory(DocumentLoaderFactory):
    def create_loader(self, folder_path: str) -> BaseLoader:
        # Define known file types and their specific loaders
        file_type_loaders = {
            ".pdf": PyPDFLoader,
            ".pptx": CustomPowerPointLoader,
            ".ppt": CustomPowerPointLoader,
            ".docx": UnstructuredWordDocumentLoader,
            ".doc": UnstructuredWordDocumentLoader,
            ".txt": TextLoader,
            ".csv": CSVLoader,
            ".json": JSONLoader,
            ".html": BSHTMLLoader,
            ".md": UnstructuredMarkdownLoader,
            ".xlsx": UnstructuredExcelLoader,
            ".xls": UnstructuredExcelLoader,
            # Programming languages
            ".py": PythonLoader,
            ".js": TextLoader,
            ".java": TextLoader,
            ".cpp": TextLoader,
            ".c": TextLoader,
            ".h": TextLoader,
            ".cs": TextLoader,
            ".php": TextLoader,
            ".rb": TextLoader,
            ".go": TextLoader,
            ".rs": TextLoader,
            ".swift": TextLoader,
            ".kt": TextLoader,
            # Web files
            ".html": TextLoader,
            ".css": TextLoader,
            ".xml": TextLoader,
            ".yaml": TextLoader,
            ".yml": TextLoader,
            # Config files
            ".ini": TextLoader,
            ".conf": TextLoader,
            ".env": TextLoader,
            ".toml": TextLoader,
            # Shell scripts
            ".sh": TextLoader,
            ".bash": TextLoader,
            ".zsh": TextLoader,
            ".fish": TextLoader,
        }
        # File extensions to ignore
        ignore_extensions = {
            ".pyc",
            ".pyo",
            ".pyd",  # Python compiled files
            ".class",
            ".jar",  # Java compiled files
            ".o",
            ".obj",
            ".exe",  # Compiled binaries
            ".dll",
            ".so",
            ".dylib",  # Libraries
            ".git",
            ".svn",  # Version control
            ".DS_Store",  # System files
        }

        loaders = []
        file_metadata = {}

        # Walk through directory and process all files
        for root, _, files in os.walk(folder_path):
            # Skip hidden directories and their contents
            if any(part.startswith(".") for part in root.split(os.sep)):
                continue

            for file in files:
                file_path = os.path.join(root, file)
                _, extension = os.path.splitext(file)
                extension = extension.lower()

                # Skip files with ignored extensions
                if extension in ignore_extensions:
                    continue

                # Get file metadata - name and modification time
                mod_time = os.path.getmtime(file_path)
                mod_date = datetime.datetime.fromtimestamp(mod_time).strftime(
                    "%Y-%m-%d"
                )

                # Store metadata for this file to be applied after loading
                file_metadata[file_path] = {
                    "filename": file,  # Just the filename without path
                    "last_modified": mod_date,
                }

                try:
                    # Use specific loader if available
                    if extension in file_type_loaders:
                        loader = file_type_loaders[extension](file_path)
                        loaders.append(loader)
                    else:
                        # Try UnstructuredFileLoader for unknown file types
                        print(f"Attempting to parse unknown file type: {file}")
                        loader = UnstructuredFileLoader(
                            file_path, mode="elements", strategy="fast"
                        )
                        loaders.append(loader)
                except Exception as e:
                    print(f"Error creating loader for {file}: {e}")
                    continue

        return MetadataEnhancedLoader(CombinedLoader(loaders), file_metadata)

    def _create_pptx_loader(self, file_path: str):
        """Create a PowerPoint loader with proper error handling"""
        try:
            return UnstructuredPowerPointLoader(
                file_path,
                mode="elements",
                strategy="fast",
                skip_download=True,  # Skip downloading models
            )
        except Exception as e:
            logging.warning(f"Failed to create PowerPoint loader: {e}")
            # Fallback to basic file loader
            return UnstructuredFileLoader(file_path, mode="elements", strategy="fast")


class CombinedLoader(BaseLoader):

    def __init__(self, loaders: List[BaseLoader]):
        self.loaders = loaders

    def load(self):
        documents = []
        for loader in self.loaders:
            try:
                docs = loader.load()
                if isinstance(docs, list):
                    documents.extend(docs)
                else:
                    documents.append(docs)
                print(f"Successfully loaded document with {type(loader).__name__}")
            except Exception as e:
                print(f"Error loading documents with {type(loader).__name__}: {e}")
                continue
        return documents

    def __str__(self):
        return f"CombinedLoader with {len(self.loaders)} loaders"


class CustomPowerPointLoader(BaseLoader):
    def __init__(self, file_path: str):
        self.file_path = file_path

    def _extract_text_from_shape(self, shape) -> Dict[str, str]:
        """Extract text and determine if it's a header based on properties"""
        if not hasattr(shape, "text") or not shape.text.strip():
            return None

        text = shape.text.strip()
        is_header = False

        try:
            # Check if shape is a title placeholder
            if shape.is_placeholder and shape.placeholder_format.type in (
                1,
                2,
            ):  # 1: Title, 2: Center Title
                is_header = True
            # else:
            #     # Check text properties for header-like characteristics
            #     for paragraph in shape.text_frame.paragraphs:
            #         if paragraph.text.strip():
            #             font = paragraph.runs[0].font if paragraph.runs else None
            #             if font and font.size:
            #                 # Convert font size from EMU to points (1 point = 12700 EMU)
            #                 size_pt = font.size / 12700
            #                 is_header = size_pt > 18  # Assume larger text is header
            #                 break

        except Exception as e:
            print(f"Error processing shape: {e}")

        return {"text": text, "is_header": is_header}

    def _process_slide(self, slide, slide_number: int) -> Tuple[List[str], List[str]]:
        """Process a single slide and return headers and contents separately"""
        headers = []
        contents = []

        # Process shapes in order of their position (top to bottom)
        shapes = [(shape, shape.top) for shape in slide.shapes if hasattr(shape, "top")]
        shapes.sort(key=lambda x: x[1])  # Sort by vertical position

        for shape, _ in shapes:
            text_info = self._extract_text_from_shape(shape)
            if text_info:
                if text_info["is_header"]:
                    headers.append(text_info["text"])
                else:
                    contents.append(text_info["text"])

        return headers, contents

    def load(self) -> List[Document]:
        try:
            prs = Presentation(self.file_path)
            documents = []

            for slide_number, slide in enumerate(prs.slides, 1):
                headers, contents = self._process_slide(slide, slide_number)

                # Format the document with headers first, then contents
                # This way, get_titles() will pick up the header as it takes the first line
                page_content = ""

                # Add headers first (this will be picked up by get_titles)
                if headers:
                    page_content += "\n".join(headers)
                else:
                    page_content += (
                        f"Slide {slide_number}"  # Fallback title if no headers
                    )

                # Add contents after headers
                if contents:
                    page_content += "\n\n" + "\n".join(contents)

                documents.append(
                    Document(
                        page_content=page_content,
                        metadata={
                            "source": self.file_path,
                            "type": "powerpoint",
                            "slide_number": slide_number,
                            "total_slides": len(prs.slides),
                            "filename": os.path.basename(self.file_path),
                        },
                    )
                )

        except Exception as e:
            print(f"Error processing PowerPoint file {self.file_path}: {e}")
            return []

        return documents

    def __str__(self):
        return f"CustomPowerPointLoader(file_path='{self.file_path}')"


class MetadataEnhancedLoader(BaseLoader):
    """A loader that adds additional metadata to documents after loading."""

    def __init__(self, loader: BaseLoader, file_metadata: Dict[str, Dict]):
        self.loader = loader
        self.file_metadata = file_metadata

    def load(self) -> List[Document]:
        documents = self.loader.load()

        for doc in documents:
            source = doc.metadata.get("source")
            if source and source in self.file_metadata:
                # Add metadata from our dictionary
                doc.metadata.update(self.file_metadata[source])

                # Also ensure filename is added (in case the source key differs)
                if "filename" not in doc.metadata:
                    doc.metadata["filename"] = os.path.basename(source)

        return documents
